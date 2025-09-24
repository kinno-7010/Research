#!/usr/bin/env python3
"""
豪雨と化学 - 中学生向け授業プレゼンテーション
PowerPointファイル生成スクリプト

このスクリプトは、豪雨と化学の関係を学ぶ教育的なプレゼンテーションを
PowerPoint形式で生成します。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

# プレゼンテーション設定
PRESENTATION_TITLE = "豪雨と化学 - 災害から学ぶ科学の力"
OUTPUT_FILENAME = f"rainfall_chemistry_presentation_{datetime.now().strftime('%Y%m%d')}.pptx"

# カラーパレット（統一感のあるデザインのため）
COLORS = {
    'primary_blue': RGBColor(52, 152, 219),    # #3498db
    'secondary_red': RGBColor(231, 76, 60),    # #e74c3c
    'dark_gray': RGBColor(44, 62, 80),         # #2c3e50
    'light_gray': RGBColor(236, 240, 241),     # #ecf0f1
    'orange': RGBColor(243, 156, 18),          # #f39c12
    'purple': RGBColor(155, 89, 182),          # #9b59b6
    'green': RGBColor(26, 188, 156),           # #1abc9c
}

def create_title_slide(prs):
    """タイトルスライドの作成"""
    slide_layout = prs.slide_layouts[0]  # タイトルスライドレイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = PRESENTATION_TITLE
    subtitle.text = "中学生向け授業用プレゼンテーション\n\n災害を科学の視点で理解し、\n未来を守る力を育む"
    
    # タイトルのフォーマット
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['dark_gray']
    
    return slide

def add_content_slide(prs, slide_number, title_text, bullet_points, visual_description):
    """コンテンツスライドの共通テンプレート"""
    # 2列レイアウトを使用
    slide_layout = prs.slide_layouts[3]  # 2列コンテンツレイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトル設定
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['dark_gray']
    
    # 左側のテキストコンテンツ
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.clear()  # 既存のテキストをクリア
    
    # 箇条書きの追加
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        
        # 重要な部分を強調
        if "科学の力" in point or "日常化した災害" in point:
            p.font.bold = True
            p.font.color.rgb = COLORS['secondary_red']
    
    # 右側のビジュアル説明（テキストボックスとして）
    right_content = slide.placeholders[2]
    right_content.text = visual_description
    
    # スライド番号の追加
    add_slide_number(slide, slide_number)
    
    return slide

def add_slide_number(slide, number):
    """スライド番号を追加"""
    left = Inches(8.5)
    top = Inches(6.8)
    width = Inches(1.5)
    height = Inches(0.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = f"{number} / 10"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLORS['dark_gray']

def create_comparison_table_slide(prs):
    """比較表スライド（スライド3）の作成"""
    slide_layout = prs.slide_layouts[5]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトル追加
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "▶ 実は、中学校の理科で習ったあれこれ"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['dark_gray']
    
    # 左側の説明
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(4), Inches(3)
    )
    left_frame = left_box.text_frame
    left_frame.text = "• 雲の発生＝飽和水蒸気量と凝結\n\n• 豪雨後の水質変化＝pH・イオン\n\n• 洪水の濁水＝懸濁物質・汚染物"
    for paragraph in left_frame.paragraphs:
        paragraph.font.size = Pt(20)
    
    # 比較表の作成
    # 左側の表（中学理科の内容）
    left_table = slide.shapes.add_table(4, 1, Inches(5), Inches(2), Inches(2), Inches(3)).table
    left_table.cell(0, 0).text = "中学理科の内容"
    left_table.cell(1, 0).text = "水の三態変化"
    left_table.cell(2, 0).text = "酸性・アルカリ性"
    left_table.cell(3, 0).text = "イオンの性質"
    
    # 右側の表（豪雨との関係）
    right_table = slide.shapes.add_table(4, 1, Inches(7.5), Inches(2), Inches(2), Inches(3)).table
    right_table.cell(0, 0).text = "豪雨との関係"
    right_table.cell(1, 0).text = "雲の形成メカニズム"
    right_table.cell(2, 0).text = "酸性雨の影響"
    right_table.cell(3, 0).text = "水質汚染の指標"
    
    # 表のスタイル設定
    for table in [left_table, right_table]:
        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(16)
                if table.rows[0] == row:  # ヘッダー行
                    cell.text_frame.paragraphs[0].font.bold = True
    
    add_slide_number(slide, 3)
    return slide

def create_theme_selection_slide(prs):
    """テーマ選択スライド（スライド7）の作成"""
    slide_layout = prs.slide_layouts[5]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "▶ 君なら、どんなつながりを見つける？"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['dark_gray']
    
    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(1.5), Inches(6), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "以下のテーマから1つを選んで考えてみよう："
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    
    # テーマカードの配置（2行3列）
    themes = [
        ("雨水と水質", "pHの変化を調べる"),
        ("浸水と材料", "防水技術の化学"),
        ("においと化学", "化学物質の特定"),
        ("川の色と成分", "濁度と溶存物質"),
        ("土のうと化学", "吸水ポリマーの仕組み"),
        ("自由テーマ", "君のアイデア！")
    ]
    
    card_width = Inches(2.8)
    card_height = Inches(1.8)
    start_x = Inches(0.7)
    start_y = Inches(2.5)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    
    for i, (theme_title, theme_desc) in enumerate(themes):
        row = i // 3
        col = i % 3
        
        left = start_x + col * (card_width + gap_x)
        top = start_y + row * (card_height + gap_y)
        
        # カード背景（色付き図形）
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['purple']
        card.line.color.rgb = COLORS['purple']
        
        # テキストの追加
        text_frame = card.text_frame
        text_frame.clear()
        text_frame.margin_top = Inches(0.3)
        
        # タイトル
        p1 = text_frame.paragraphs[0]
        p1.text = theme_title
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        
        # 説明
        p2 = text_frame.add_paragraph()
        p2.text = theme_desc
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(255, 255, 255)
    
    add_slide_number(slide, 7)
    return slide

def create_summary_slide(prs):
    """まとめスライド（スライド10）の作成"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "▶ 科学は、災害から社会を守る力になる"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['dark_gray']
    
    # メインメッセージ
    main_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(2.5)
    )
    main_frame = main_box.text_frame
    main_frame.text = """• 「雨」は身近な自然現象。でも「災害」にもなる

• その現場には、たくさんの化学の知識がある

• "自分ごと"にして考える力が、未来をつくる"""
    
    for paragraph in main_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.line_spacing = 1.5
    
    # キーワードのマインドマップ（簡略版）
    # 中心の円
    center_x = Inches(5)
    center_y = Inches(5.5)
    
    center_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x - Inches(1), center_y - Inches(0.5),
        Inches(2), Inches(1)
    )
    center_shape.fill.solid()
    center_shape.fill.fore_color.rgb = COLORS['secondary_red']
    center_shape.text_frame.text = "豪雨×化学"
    center_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    center_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    center_shape.text_frame.paragraphs[0].font.bold = True
    
    # 4つのキーワード
    keywords = [
        ("科学", COLORS['primary_blue'], -2, -1),
        ("災害", COLORS['purple'], 2, -1),
        ("日常", COLORS['green'], -2, 1),
        ("未来", COLORS['orange'], 2, 1)
    ]
    
    for keyword, color, x_offset, y_offset in keywords:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            center_x + Inches(x_offset) - Inches(0.7),
            center_y + Inches(y_offset) - Inches(0.35),
            Inches(1.4), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.text_frame.text = keyword
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        shape.text_frame.paragraphs[0].font.bold = True
    
    add_slide_number(slide, 10)
    return slide

def main():
    """メイン処理"""
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドサイズの設定（16:9）
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 各スライドの作成
    print("プレゼンテーションを作成中...")
    
    # タイトルスライド
    create_title_slide(prs)
    
    # スライド1: 導入
    add_content_slide(
        prs, 1,
        "▶ 私たちはなぜ「豪雨」を学ぶのか？",
        [
            "毎年のように発生する大雨・洪水",
            "命・暮らし・経済に大きな影響",
            "災害に立ち向かうには「科学の力」が必要"
        ],
        "【視覚要素】\n・熱海土石流災害\n  死者・行方不明者: 28名\n  被害額: 約600億円\n\n・熊本豪雨災害\n  死者: 67名\n  被害額: 約5,800億円\n\n近年の豪雨被害は増加傾向"
    )
    
    # スライド2: 気候変動
    add_content_slide(
        prs, 2,
        "▶ 雨の降り方が変わってきた？",
        [
            "地球温暖化 → 空気中の水蒸気量UP",
            "線状降水帯・ゲリラ豪雨の頻度増加",
            "豪雨は「日常化した災害」へ"
        ],
        "【視覚要素】\n・気温と降水量の年変化グラフ\n  →上昇傾向が明確\n\n・線状降水帯の模式図\n  →同じ場所に長時間\n  →強い雨が継続\n\n気象庁データより"
    )
    
    # スライド3: 比較表
    create_comparison_table_slide(prs)
    
    # スライド4: 水質問題
    add_content_slide(
        prs, 4,
        "▶ 洪水の後、水道水が変なにおいがする？",
        [
            "浸水 → 下水や工場排水が逆流",
            "雨水と混じってアンモニア・有機物発生",
            "においの正体＝化学物質のサイン"
        ],
        "【においの化学物質】\n\nNH₃（アンモニア）\n→ 刺激臭\n\nH₂S（硫化水素）\n→ 腐卵臭\n\n化学物質の検出は\n水質汚染の重要な指標"
    )
    
    # スライド5: 防災技術
    add_content_slide(
        prs, 5,
        "▶ 実は、災害現場の\"裏方\"は化学",
        [
            "吸水ポリマーでできた「簡易土のう」",
            "河川水の重金属分析（鉛、ヒ素など）",
            "災害時の水質検査キット＝化学反応の応用"
        ],
        "【吸水ポリマーの威力】\n\n吸水前：\n軽量・コンパクト\n（約400g）\n\n↓ 水を吸収\n\n吸水後：\n約20kgの土のうに！\n\n保管・運搬が容易で\n緊急時に即対応可能")
    
    # スライド6: 環境保護
    add_content_slide(
        prs, 6,
        "▶ 水だけじゃない、空気も土も化学で守る",
        [
            "雨水は空気の化学物質を溶かし込む（酸性雨）",
            "土壌中に流れ込んだ汚染物質の拡散を防ぐには？",
            "除染・吸着材・中和剤 → すべて化学の知見"
        ],
        "【日本の酸性雨状況】\n\n都市部：pH 4.5以下\n郊外：pH 5.0-5.5\n\n酸性雨の影響：\n・建物の劣化\n・土壌の酸性化\n・植物への影響\n\n化学的対策が不可欠"
    )
    
    # スライド7: テーマ選択
    create_theme_selection_slide(prs)
    
    # スライド8: ワークタイム
    add_content_slide(
        prs, 8,
        "▶ 実際に「科学の目」で考えてみよう！",
        [
            "どこに化学が関係していそう？",
            "使われている化学物質や原理は？",
            "自分ならどんな工夫をする？"
        ],
        "【ワークシート】\n\n1. 選んだテーマ：______\n\n2. 化学との関係：\n___________________\n\n3. 実験アイデア：\n___________________\n\n考え方のヒント：\n身の回りの「当たり前」を\n疑ってみよう！"
    )
    
    # スライド9: 発表
    add_content_slide(
        prs, 9,
        "▶ 君の「豪雨×化学」の発見を共有しよう！",
        [
            "「においの原因をアンモニアと仮定」",
            "「撥水加工の化学構造を調べた」",
            "「学校の雨どいの水を測ってみたい」"
        ],
        "【発表のポイント】\n\n✓ 化学の視点で説明\n\n✓ 具体的な物質名や反応\n\n✓ 実験可能なアイデア\n\n✓ 日常生活との関連\n\n相手に伝わりやすく\n科学的根拠を明確に！"
    )
    
    # スライド10: まとめ
    create_summary_slide(prs)
    
    # ファイルの保存
    output_path = os.path.join(os.getcwd(), OUTPUT_FILENAME)
    prs.save(output_path)
    print(f"\nプレゼンテーションファイルを作成しました: {output_path}")
    print(f"ファイルサイズ: {os.path.getsize(output_path) / 1024:.1f} KB")
    print("\n授業でのご活用をお祈りしています！")

if __name__ == "__main__":
    main()